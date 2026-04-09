"""
Multimodal ingestion pipeline.

Drop any supported file type into the data folder and the pipeline will:

  1. Discover all files (text, PDF, image, audio, video, Office docs).
  2. Extract text content for text-based files (used as a human-readable
     document store).
  3. Build multimodal embedding inputs for every file so that Qwen3-VL-Embedding
     can encode each item natively — images and videos are passed as raw pixels,
     not converted to text first.

Embedding model
---------------
  Qwen/Qwen3-VL-Embedding-2B   (default, ~5 GB VRAM / CPU-offloadable)
  Qwen/Qwen3-VL-Embedding-8B   (better quality, ~18 GB VRAM)

  Both run fully locally via HuggingFace transformers — no API key required.
  On first use the weights are downloaded to ~/.cache/huggingface.
  You can also pass a local directory path as model_id.

Reranker model
--------------
  Qwen/Qwen3-VL-Reranker-2B
  Qwen/Qwen3-VL-Reranker-8B

Supported file types
--------------------
  Text      : .txt  .md  .rst  .csv  .json  .xml  .html  .htm
  PDF       : .pdf   (pdfminer.six > pypdf fallback for text extraction)
  Images    : .jpg  .jpeg  .png  .gif  .bmp  .tiff  .webp
  Audio     : .mp3  .wav  .ogg  .flac  .m4a         (Whisper transcription)
  Video     : .mp4  .avi  .mov  .mkv  .webm          (native to Qwen3-VL OR
                                                       Whisper audio-track fallback)
  Office    : .docx  .pptx  .xlsx

Key design decisions
--------------------
- Text files, PDFs, and Office docs: text is extracted then passed as {"text": ...}
  to Qwen3-VL-Embedding.
- Images: passed as {"image": path} — the model encodes pixels natively.
- Videos: passed as {"video": path} — the model samples frames natively.
- Audio: transcribed to text with Whisper (audio has no visual tokens).
- Every item therefore gets a true multimodal embedding in a shared vector space.
- All dependencies are optional with graceful keyword-search fallback.

Dependencies (all optional — graceful fallback when missing)
------------------------------------------------------------
  Core (required for Qwen3-VL):
    transformers>=4.57.3
    torch>=2.0
    qwen-vl-utils>=0.0.14
    accelerate>=1.12.0
    scipy>=1.16.0
    pillow
    faiss-cpu          (or faiss-gpu)

  PDF text:
    pdfminer.six  OR  pypdf

  Audio / video audio track:
    openai-whisper     (pip install openai-whisper)
    ffmpeg on PATH     (for non-WAV formats)

  Office files:
    python-docx        (pip install python-docx)
    openpyxl           (pip install openpyxl)
    python-pptx        (pip install python-pptx)
"""

from __future__ import annotations

import re
import sys
from pathlib import Path
from typing import Any


# ---------------------------------------------------------------------------
# Supported extensions
# ---------------------------------------------------------------------------

_TEXT_EXTS  = {".txt", ".md", ".rst", ".csv", ".json", ".xml", ".html", ".htm"}
_PDF_EXTS   = {".pdf"}
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".tif", ".webp"}
_AUDIO_EXTS = {".mp3", ".wav", ".ogg", ".flac", ".m4a", ".aac"}
_VIDEO_EXTS = {".mp4", ".avi", ".mov", ".mkv", ".webm"}
_OFFICE_EXTS = {".docx", ".pptx", ".xlsx"}

ALL_SUPPORTED_EXTS = (
    _TEXT_EXTS | _PDF_EXTS | _IMAGE_EXTS | _AUDIO_EXTS | _VIDEO_EXTS | _OFFICE_EXTS
)

# Extension → modality tag
_MODALITY: dict[str, str] = {}
for _e in _TEXT_EXTS:   _MODALITY[_e] = "text"
for _e in _PDF_EXTS:    _MODALITY[_e] = "pdf"
for _e in _IMAGE_EXTS:  _MODALITY[_e] = "image"
for _e in _AUDIO_EXTS:  _MODALITY[_e] = "audio"
for _e in _VIDEO_EXTS:  _MODALITY[_e] = "video"
for _e in _OFFICE_EXTS: _MODALITY[_e] = "office"


# ---------------------------------------------------------------------------
# Lazy model singletons (one instance per process, shared across tools)
# ---------------------------------------------------------------------------

_embedder_instance: Any = None   # Qwen3VLEmbedder
_reranker_instance: Any = None   # Qwen3VLReranker
_whisper_instance:  Any = None


def get_embedder(model_id: str = "Qwen/Qwen3-VL-Embedding-2B", **kwargs: Any) -> Any:
    """
    Return a cached Qwen3VLEmbedder instance (loaded once per process).

    Parameters
    ----------
    model_id : str
        HuggingFace model id OR local directory path.
        Examples:
          "Qwen/Qwen3-VL-Embedding-2B"   – auto-downloads to HF cache
          "Qwen/Qwen3-VL-Embedding-8B"   – larger, better quality
          "/local/path/Qwen3-VL-Embedding-2B"  – fully offline
    **kwargs :
        Forwarded to Qwen3VLEmbedder (e.g. torch_dtype, attn_implementation).
    """
    global _embedder_instance
    if _embedder_instance is None:
        # Qwen3VLEmbedder lives in a `scripts/` file bundled inside the model
        # repo on HuggingFace — it is NOT an installable package.
        # We download (or reuse from cache) that script via hf_hub_download,
        # then import it dynamically.
        import importlib.util
        from huggingface_hub import hf_hub_download  # type: ignore[import]

        # model_id may be a local path already containing the script
        scripts_local = Path(model_id) / "scripts" / "qwen3_vl_embedding.py"
        if scripts_local.exists():
            script_path = str(scripts_local)
        else:
            # derive the repo_id: strip any local-path prefix, keep "Org/Model"
            repo_id = model_id if "/" in model_id else f"Qwen/{model_id}"
            script_path = hf_hub_download(
                repo_id=repo_id,
                filename="scripts/qwen3_vl_embedding.py",
            )

        spec = importlib.util.spec_from_file_location("qwen3_vl_embedding", script_path)
        mod = importlib.util.module_from_spec(spec)  # type: ignore[arg-type]
        spec.loader.exec_module(mod)  # type: ignore[union-attr]
        Qwen3VLEmbedder = mod.Qwen3VLEmbedder
        _embedder_instance = Qwen3VLEmbedder(model_name_or_path=model_id, **kwargs)
    return _embedder_instance


def get_reranker(model_id: str = "Qwen/Qwen3-VL-Reranker-2B", **kwargs: Any) -> Any:
    """
    Return a cached Qwen3VLReranker instance (loaded once per process).

    Parameters
    ----------
    model_id : str
        HuggingFace model id OR local directory path.
    """
    global _reranker_instance
    if _reranker_instance is None:
        # Same pattern as get_embedder() — script lives inside the model repo.
        import importlib.util
        from huggingface_hub import hf_hub_download  # type: ignore[import]

        scripts_local = Path(model_id) / "scripts" / "qwen3_vl_reranker.py"
        if scripts_local.exists():
            script_path = str(scripts_local)
        else:
            repo_id = model_id if "/" in model_id else f"Qwen/{model_id}"
            script_path = hf_hub_download(
                repo_id=repo_id,
                filename="scripts/qwen3_vl_reranker.py",
            )

        spec = importlib.util.spec_from_file_location("qwen3_vl_reranker", script_path)
        mod = importlib.util.module_from_spec(spec)  # type: ignore[arg-type]
        spec.loader.exec_module(mod)  # type: ignore[union-attr]
        Qwen3VLReranker = mod.Qwen3VLReranker
        _reranker_instance = Qwen3VLReranker(model_name_or_path=model_id, **kwargs)
    return _reranker_instance


def get_whisper(model_size: str = "base") -> Any:
    """Return a cached Whisper model."""
    global _whisper_instance
    if _whisper_instance is None:
        import whisper  # type: ignore[import]
        _whisper_instance = whisper.load_model(model_size)
    return _whisper_instance


# ---------------------------------------------------------------------------
# Text extractors (produce human-readable strings for entity extraction
# and the document store; separately, we build embedding inputs below)
# ---------------------------------------------------------------------------

def _read_text_file(path: Path) -> str:
    """Read plain-text with automatic encoding detection."""
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return path.read_text(errors="replace")


def _extract_pdf_text(path: Path) -> str:
    """Extract text from PDF: pdfminer.six > pypdf fallback."""
    try:
        from pdfminer.high_level import extract_text  # type: ignore[import]
        text = extract_text(str(path))
        if text and text.strip():
            return text
    except (ImportError, Exception):
        pass

    try:
        from pypdf import PdfReader  # type: ignore[import]
        reader = PdfReader(str(path))
        return "\n\n".join(p.extract_text() or "" for p in reader.pages)
    except (ImportError, Exception):
        pass

    return ""


def _render_pdf_pages_with_text(
    path: Path,
    dpi: int = 150,
    max_pages: int = 0,
) -> list[tuple[Path, str]]:
    """
    Render each page of a PDF to a temporary PNG **and** extract the page's
    text in a single PyMuPDF pass.

    Returns a list of ``(png_path, page_text)`` tuples — one per page.
    Both the image path and the extracted text are always present; if a page
    has no selectable text (e.g. a scanned figure), ``page_text`` is ``""``.

    Using PyMuPDF for both operations avoids opening the file twice and
    guarantees that the text stored in the chunk index exactly matches the
    layout of the rendered image, enabling downstream agents to reference
    text from a visually-retrieved page.

    Falls back to an empty list when PyMuPDF is not installed — the pipeline
    will still embed the full-document text chunks.

    Parameters
    ----------
    path : Path
        PDF file to render.
    dpi : int
        Render resolution.  150 dpi balances quality vs. VRAM/memory.
        Use 120 for very large PDFs; 200 for dense figures.
    max_pages : int
        Cap on pages rendered per PDF.  0 = no cap (render all pages).

    Returns
    -------
    list[tuple[Path, str]]
        ``(png_path, page_text)`` pairs (one per page).
    """
    try:
        import fitz  # type: ignore[import]  # PyMuPDF
        import tempfile

        doc = fitz.open(str(path))
        mat = fitz.Matrix(dpi / 72, dpi / 72)  # 72 pt/inch → dpi
        tmp_dir = Path(tempfile.gettempdir()) / "dlc_pdf_pages"
        tmp_dir.mkdir(exist_ok=True)

        results: list[tuple[Path, str]] = []
        n_pages = len(doc)
        limit = n_pages if max_pages <= 0 else min(max_pages, n_pages)

        for i in range(limit):
            page = doc[i]
            # Render to PNG
            pix = page.get_pixmap(matrix=mat, alpha=False)
            out_path = tmp_dir / f"{path.stem}_p{i:04d}.png"
            pix.save(str(out_path))
            # Extract selectable text from the same page object
            page_text = page.get_text("text").strip()
            results.append((out_path, page_text))

        doc.close()
        return results

    except ImportError:
        return []
    except Exception:
        return []


# Thin backward-compatible wrapper (returns only paths — used by nothing
# in the codebase anymore, kept so external callers don't break).
def _render_pdf_pages(
    path: Path,
    dpi: int = 150,
    max_pages: int = 0,
) -> list[Path]:
    """Backward-compatible wrapper — prefer ``_render_pdf_pages_with_text``."""
    return [p for p, _ in _render_pdf_pages_with_text(path, dpi, max_pages)]


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document  # type: ignore[import]
        doc = Document(str(path))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return f"[DOCX: {path.name} — install python-docx]"
    except Exception as exc:
        return f"[DOCX: {path.name} — {exc}]"


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore[import]
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        return f"[PPTX: {path.name} — install python-pptx]"
    except Exception as exc:
        return f"[PPTX: {path.name} — {exc}]"


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl  # type: ignore[import]
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = ["\t".join(str(c) for c in row if c is not None)
                    for row in ws.iter_rows(values_only=True)]
            rows = [r for r in rows if r.strip()]
            if rows:
                sheets.append(f"[Sheet: {name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)
    except ImportError:
        return f"[XLSX: {path.name} — install openpyxl]"
    except Exception as exc:
        return f"[XLSX: {path.name} — {exc}]"


def _transcribe_audio(path: Path, whisper_model_size: str = "base") -> str:
    """Transcribe audio with Whisper (runs locally, no API)."""
    try:
        model = get_whisper(whisper_model_size)
        result = model.transcribe(str(path), fp16=False)
        transcript = result.get("text", "").strip()
        return f"[AUDIO TRANSCRIPT]\n{transcript}" if transcript else ""
    except ImportError:
        return f"[AUDIO: {path.name} — install openai-whisper to transcribe]"
    except Exception as exc:
        return f"[AUDIO: {path.name} — transcription failed: {exc}]"


# ---------------------------------------------------------------------------
# Embedding-input builders
# Each function returns a dict accepted by Qwen3VLEmbedder.process([...])
# ---------------------------------------------------------------------------

def _embed_input_for_text(text: str, instruction: str | None = None) -> dict[str, Any]:
    """Text → embedding input dict."""
    d: dict[str, Any] = {"text": text}
    if instruction:
        d["instruction"] = instruction
    return d


def _embed_input_for_image(path: Path, instruction: str | None = None) -> dict[str, Any]:
    """Image → embedding input dict (pixels encoded natively by the model)."""
    d: dict[str, Any] = {"image": str(path.resolve())}
    if instruction:
        d["instruction"] = instruction
    return d


def _embed_input_for_video(
    path: Path,
    fps: float = 1.0,
    max_frames: int = 64,
    instruction: str | None = None,
) -> dict[str, Any]:
    """Video → embedding input dict (frames sampled + encoded natively)."""
    d: dict[str, Any] = {"video": str(path.resolve()), "fps": fps, "max_frames": max_frames}
    if instruction:
        d["instruction"] = instruction
    return d


# ---------------------------------------------------------------------------
# Chunking helper (for the text document store)
# ---------------------------------------------------------------------------

def _chunk_text(text: str, chunk_size: int = 1500, overlap: int = 200) -> list[str]:
    """Split text into overlapping chunks on sentence boundaries."""
    if len(text) <= chunk_size:
        return [text]

    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks: list[str] = []
    current: list[str] = []
    current_len = 0

    for sent in sentences:
        if current_len + len(sent) > chunk_size and current:
            chunks.append(" ".join(current))
            # Keep overlap
            overlap_buf: list[str] = []
            ol = 0
            for s in reversed(current):
                if ol + len(s) > overlap:
                    break
                overlap_buf.insert(0, s)
                ol += len(s)
            current = overlap_buf
            current_len = ol
        current.append(sent)
        current_len += len(sent)

    if current:
        chunks.append(" ".join(current))

    return chunks or [text]


# ---------------------------------------------------------------------------
# Main pipeline class
# ---------------------------------------------------------------------------

class MultimodalIngestionPipeline:
    """
    Ingest any mix of files from a folder and produce:

      - ``documents``   dict[str, str]            doc_id → extracted text
      - ``chunk_index`` list[dict]                per-chunk metadata
      - ``embed_inputs`` list[dict]               Qwen3VL-ready input per chunk
      - (optionally) a FAISS index via ``build_faiss_index()``

    Usage
    -----
    ::

        pipeline = MultimodalIngestionPipeline(
            data_dir="my_data/",
            embedding_model_id="Qwen/Qwen3-VL-Embedding-2B",
        )
        documents = pipeline.run()
        faiss_index = pipeline.build_faiss_index()

    All model weights are loaded locally from HuggingFace cache or a local
    directory.  No API keys are required.

    Parameters
    ----------
    data_dir : str | Path
        Folder to scan for supported files.
    recursive : bool
        Also scan sub-folders.
    min_text_length : int
        Skip files that produce fewer than this many characters of text.
    embedding_model_id : str
        Local path or HuggingFace id for Qwen3-VL-Embedding.
    embedding_kwargs : dict
        Extra kwargs forwarded to Qwen3VLEmbedder (e.g. ``torch_dtype``,
        ``attn_implementation``).
    whisper_model_size : str
        Whisper checkpoint size for audio transcription.
    video_fps : float
        Frame sampling rate for video inputs.
    video_max_frames : int
        Maximum frames to sample per video.
    embed_instruction : str | None
        Instruction prefix for document embeddings (improves recall ~1-5 %).
    chunk_size : int
        Target character length per text chunk.
    chunk_overlap : int
        Character overlap between consecutive chunks.
    verbose : bool
        Print progress to stdout.
    """

    _DEFAULT_INSTRUCTION = (
        "Represent this document for retrieval in a knowledge discovery system. "
        "Capture entities, relationships, and domain-specific facts."
    )

    def __init__(
        self,
        data_dir: str | Path,
        min_text_length: int = 30,
        embedding_model_id: str = "Qwen/Qwen3-VL-Embedding-2B",
        embedding_kwargs: dict[str, Any] | None = None,
        whisper_model_size: str = "base",
        video_fps: float = 1.0,
        video_max_frames: int = 64,
        embed_instruction: str | None = None,
        chunk_size: int = 1500,
        chunk_overlap: int = 200,
        pdf_dpi: int = 150,
        pdf_max_pages: int = 0,
        verbose: bool = True,
    ) -> None:
        self.data_dir = Path(data_dir)
        self.min_text_length = min_text_length
        self.embedding_model_id = embedding_model_id
        self.embedding_kwargs: dict[str, Any] = embedding_kwargs or {}
        self.whisper_model_size = whisper_model_size
        self.video_fps = video_fps
        self.video_max_frames = video_max_frames
        self.embed_instruction = embed_instruction or self._DEFAULT_INSTRUCTION
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.pdf_dpi = pdf_dpi
        self.pdf_max_pages = pdf_max_pages
        self.verbose = verbose

        # Populated after run()
        self.documents:    dict[str, str]       = {}
        self.chunk_index:  list[dict[str, Any]] = []  # metadata per embeddable item
        self.embed_inputs: list[dict[str, Any]] = []  # Qwen3VL input dicts
        self.embeddings:   Any                  = None  # np.ndarray (n, dim)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def run(self) -> dict[str, str]:
        """
        Discover, extract, embed, and index all files.

        Returns
        -------
        documents : dict[str, str]
        """
        paths = self._discover_files()
        if not paths:
            self._log(
                f"WARNING: No supported files found in '{self.data_dir}'.\n"
                f"Supported: {', '.join(sorted(ALL_SUPPORTED_EXTS))}"
            )
            return {}

        self._log(f"Found {len(paths)} file(s) in '{self.data_dir}'")

        # ---- Phase 1: per-file extraction --------------------------------
        for path in paths:
            doc_id   = self._doc_id(path)
            modality = _MODALITY.get(path.suffix.lower(), "text")
            self._log(f"  [{modality:6s}] {path.name}  →  '{doc_id}'")

            try:
                text, embed_items = self._process_file(path, modality)
            except Exception as exc:
                self._log(f"    SKIP — {exc}")
                continue

            if not text.strip() and not embed_items:
                self._log("    SKIP — no content extracted")
                continue

            # Store text representation (may be empty for pure image/video)
            if text.strip():
                if len(text.strip()) < self.min_text_length:
                    self._log(f"    SKIP — too little text ({len(text.strip())} chars)")
                    continue
                self.documents[doc_id] = text
                self._log(f"    {len(text):,} chars extracted")
            else:
                # Image/video with no text — still store a placeholder
                self.documents[doc_id] = f"[{modality.upper()}: {path.name}]"

            # Register embedding inputs
            for item in embed_items:
                # For text chunks: use the chunk text directly.
                # For PDF page images: fall back to "page_text" (the text
                # extracted from that specific page by PyMuPDF) so that
                # downstream agents can read meaningful content even when the
                # retrieval was triggered by a visual/layout match.
                full_text = item.get("text") or item.get("page_text") or ""
                self.chunk_index.append({
                    "doc_id":        doc_id,
                    "chunk_id":      f"{doc_id}__item{len(self.chunk_index)}",
                    "modality":      modality,
                    "text_snippet":  full_text[:300],
                    "text":          full_text,   # full chunk text for downstream processing
                })
                self.embed_inputs.append(item)

        if not self.documents:
            self._log("WARNING: No files produced usable content.")
            return {}

        self._log(
            f"\nIngestion summary: {len(self.documents)} documents, "
            f"{len(self.embed_inputs)} embeddable items"
        )

        # ---- Phase 2: embed ---------------------------------------------
        self._build_embeddings()

        return self.documents

    def build_faiss_index(self) -> Any:
        """
        Build a FAISS IndexFlatIP from the chunk embeddings.

        Returns None if FAISS is not installed or embeddings were not built.
        """
        if self.embeddings is None or len(self.chunk_index) == 0:
            return None
        try:
            import faiss       # type: ignore[import]
            import numpy as np # type: ignore[import]

            emb = self.embeddings.astype(np.float32)
            faiss.normalize_L2(emb)
            index = faiss.IndexFlatIP(emb.shape[1])
            index.add(emb)
            return index
        except ImportError:
            return None

    # ------------------------------------------------------------------
    # File discovery
    # ------------------------------------------------------------------

    def _discover_files(self) -> list[Path]:
        if not self.data_dir.exists():
            self._log(f"ERROR: '{self.data_dir}' does not exist.")
            return []

        glob = self.data_dir.rglob
        paths: list[Path] = []
        for ext in ALL_SUPPORTED_EXTS:
            paths.extend(glob(f"*{ext}"))
            paths.extend(glob(f"*{ext.upper()}"))
        return sorted(set(paths))

    def _doc_id(self, path: Path) -> str:
        try:
            rel = path.relative_to(self.data_dir)
            return str(rel.with_suffix("")).replace("/", "__").replace("\\", "__")
        except ValueError:
            return path.stem

    # ------------------------------------------------------------------
    # Per-file processing
    # Returns (text, list_of_embed_input_dicts)
    # ------------------------------------------------------------------

    def _process_file(
        self, path: Path, modality: str
    ) -> tuple[str, list[dict[str, Any]]]:
        """
        Extract text representation and build Qwen3VL embedding input(s).

        For text-based modalities the text is chunked and each chunk becomes
        one embedding item ({"text": chunk}).

        For images and videos the raw file is passed directly so the model
        encodes pixels natively — no intermediate text conversion.

        For audio, Whisper transcribes to text which is then chunked.
        """
        instr = self.embed_instruction

        if modality == "text":
            text = _read_text_file(path)
            items = [
                _embed_input_for_text(chunk, instr)
                for chunk in _chunk_text(text, self.chunk_size, self.chunk_overlap)
            ]
            return text, items

        if modality == "pdf":
            text = _extract_pdf_text(path)
            # Text chunks — used by NER, entity extraction, and document store.
            items: list[dict[str, Any]] = [
                _embed_input_for_text(chunk, instr)
                for chunk in _chunk_text(text, self.chunk_size, self.chunk_overlap)
            ]
            # Hybrid page items: PNG image (so Qwen3-VL-Embedding encodes layout,
            # figures, and tables natively) PLUS the page's own extracted text
            # stored in the embed dict under "page_text".  This ensures that when
            # a visual chunk is retrieved by the search tool, downstream agents
            # can read the text that belongs to that specific page.
            page_pairs = _render_pdf_pages_with_text(
                path,
                dpi=self.pdf_dpi,
                max_pages=self.pdf_max_pages,
            )
            for pp, pt in page_pairs:
                # Send image AND page text together in one dict so the model
                # receives both visual and textual context in a single forward
                # pass — format_model_input() handles both keys simultaneously.
                # For scanned pages (pt == "") the text key is omitted so the
                # model still gets a pure-image embedding rather than "NULL".
                img_item = _embed_input_for_image(pp, instr)
                if pt:
                    img_item["text"] = pt    # multimodal: image + text together
                img_item["page_text"] = pt   # also stored for chunk_index metadata
                items.append(img_item)
            if page_pairs:
                self._log(
                    f"    + {len(page_pairs)} page image(s) added for visual embedding"
                )
            return text, items

        if modality == "image":
            # Qwen3-VL-Embedding takes the image directly — native pixel encoding
            items = [_embed_input_for_image(path, instr)]
            return "", items  # no text; document store gets placeholder

        if modality == "video":
            # Qwen3-VL-Embedding samples frames natively from the video file
            items = [_embed_input_for_video(
                path, fps=self.video_fps,
                max_frames=self.video_max_frames,
                instruction=instr,
            )]
            return "", items

        if modality == "audio":
            text = _transcribe_audio(path, self.whisper_model_size)
            items = [
                _embed_input_for_text(chunk, instr)
                for chunk in _chunk_text(text, self.chunk_size, self.chunk_overlap)
            ]
            return text, items

        if modality == "office":
            ext = path.suffix.lower()
            if ext == ".docx":
                text = _extract_docx(path)
            elif ext == ".pptx":
                text = _extract_pptx(path)
            elif ext == ".xlsx":
                text = _extract_xlsx(path)
            else:
                text = ""
            items = [
                _embed_input_for_text(chunk, instr)
                for chunk in _chunk_text(text, self.chunk_size, self.chunk_overlap)
            ]
            return text, items

        return "", []

    # ------------------------------------------------------------------
    # Embedding
    # ------------------------------------------------------------------

    def _build_embeddings(self) -> None:
        """Embed all items using Qwen3-VL-Embedding (local inference)."""
        if not self.embed_inputs:
            return

        self._log(
            f"\nEmbedding {len(self.embed_inputs)} items with "
            f"'{self.embedding_model_id}' (local) …"
        )

        try:
            embedder = get_embedder(self.embedding_model_id, **self.embedding_kwargs)
            import torch  # type: ignore[import]

            # Process in batches to avoid OOM
            batch_size = 8
            all_vecs = []
            for i in range(0, len(self.embed_inputs), batch_size):
                batch = self.embed_inputs[i : i + batch_size]
                vecs = embedder.process(batch)          # returns normalised tensor
                all_vecs.append(vecs.cpu().float().numpy())
                self._log(
                    f"    Embedded {min(i + batch_size, len(self.embed_inputs))}"
                    f"/{len(self.embed_inputs)}"
                )

            import numpy as np  # type: ignore[import]
            self.embeddings = np.vstack(all_vecs)
            self._log(
                f"Embedding matrix: {self.embeddings.shape}  "
                f"(dim={self.embeddings.shape[1]})"
            )

        except ImportError as exc:
            self._log(
                f"WARNING: Cannot load Qwen3-VL-Embedding — {exc}\n"
                "  Install: pip install transformers>=4.57.3 torch qwen-vl-utils accelerate\n"
                "  Falling back to keyword search."
            )
            self.embeddings = None
        except Exception as exc:
            self._log(f"WARNING: Embedding failed ({exc}). Falling back to keyword search.")
            self.embeddings = None

    # ------------------------------------------------------------------

    def _log(self, msg: str) -> None:
        if self.verbose:
            print(msg, flush=True)
